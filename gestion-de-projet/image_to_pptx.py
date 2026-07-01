#!/usr/bin/env python3
"""
image_to_pptx.py — Convertit une image design en slide PowerPoint entièrement éditable.

Claude Vision analyse le layout (positions, couleurs, typographie) et python-pptx
reconstruit chaque élément comme un objet PowerPoint natif modifiable.

Usage:
  python image_to_pptx.py <image> [sortie.pptx]
  python image_to_pptx.py <image.json>          # rejouer depuis une analyse existante

Dépendances:
  pip install anthropic python-pptx Pillow

Variable d'env requise:
  ANTHROPIC_API_KEY
"""

import sys
import base64
import json
import os
from pathlib import Path

ANALYSIS_PROMPT = """
Analyse cette image de design et extrais tous les éléments visuels pour la recréer
fidèlement en slide PowerPoint éditable.

Retourne UNIQUEMENT un objet JSON valide, sans explication, avec cette structure exacte :

{
  "slide": {
    "aspect_ratio": "16:9",
    "background_color": "#FFFFFF"
  },
  "elements": [
    {
      "id": 1,
      "type": "shape|text|shape_with_text",
      "x_pct": 0.0,
      "y_pct": 0.0,
      "w_pct": 100.0,
      "h_pct": 100.0,
      "z_order": 0,
      "shape_type": "rectangle|rounded_rectangle|ellipse|triangle|diamond",
      "fill_color": "#RRGGBB",
      "border_color": null,
      "border_width_pt": 0,
      "corner_radius_pct": 0,
      "text": "",
      "font_family": "Arial",
      "font_size_pt": 24,
      "font_bold": false,
      "font_italic": false,
      "font_color": "#000000",
      "text_align": "left|center|right",
      "vertical_align": "top|middle|bottom"
    }
  ]
}

Règles :
- x_pct/y_pct/w_pct/h_pct : pourcentages des dimensions du slide (0.0 à 100.0)
- z_order : 0 = arrière-plan, plus haut = devant
- "shape" : forme sans texte ; "text" : texte seul fond transparent ; "shape_with_text" : forme avec texte dedans
- Identifier TOUS les éléments : fonds, formes, blocs texte, icônes (formes géométriques), images (rectangle placeholder)
- Estimer les tailles de police comme elles apparaîtraient sur un slide 1920×1080
- Couleurs en hex (#RRGGBB) ou null
- Être précis sur les positions pour correspondre au layout original
"""

SHAPE_IDS = {
    "rectangle": 1,
    "rounded_rectangle": 5,
    "ellipse": 9,
    "circle": 9,
    "oval": 9,
    "triangle": 7,
    "diamond": 4,
    "pentagon": 56,
    "hexagon": 10,
    "right_arrow": 13,
    "left_arrow": 66,
    "star_5": 94,
}


def hex_to_rgb(color_str):
    if not color_str or str(color_str).lower() in ("null", "none", ""):
        return None
    try:
        from pptx.dml.color import RGBColor
        h = str(color_str).lstrip("#").strip()
        if len(h) == 3:
            h = h[0] * 2 + h[1] * 2 + h[2] * 2
        if len(h) != 6:
            return None
        return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))
    except Exception:
        return None


def resize_image_bytes(image_path: str, max_px: int = 800) -> tuple[bytes, str]:
    from PIL import Image
    import io
    img = Image.open(image_path).convert("RGB")
    img.thumbnail((max_px, max_px), Image.LANCZOS)
    buf = io.BytesIO()
    img.save(buf, format="JPEG", quality=75, optimize=True)
    return buf.getvalue(), "image/jpeg"


def analyze_image(image_path: str) -> dict:
    import anthropic
    import io
    from PIL import Image as PILImage

    api_key = os.environ.get("ANTHROPIC_API_KEY")
    if not api_key:
        # Fallback: read from scripts/ai_config.json
        config_path = Path(__file__).parent / "ai_config.json"
        if config_path.exists():
            with open(config_path) as f:
                cfg = json.load(f)
            api_key = cfg.get("anthropic_api_key", "")
    if not api_key or api_key.startswith("sk-ant-VOTRE"):
        print("Erreur : clé API Anthropic non configurée.")
        print(f"Mets ta clé dans {Path(__file__).parent / 'ai_config.json'}")
        sys.exit(1)

    print("  Analyse via Claude Vision API...")
    # Resize to reduce token usage
    img = PILImage.open(image_path).convert("RGB")
    img.thumbnail((1200, 1200), PILImage.LANCZOS)
    buf = io.BytesIO()
    img.save(buf, format="JPEG", quality=80)
    img_b64 = base64.standard_b64encode(buf.getvalue()).decode("utf-8")

    client = anthropic.Anthropic(api_key=api_key)
    message = client.messages.create(
        model="claude-opus-4-5",
        max_tokens=4096,
        messages=[{
            "role": "user",
            "content": [
                {"type": "image", "source": {"type": "base64", "media_type": "image/jpeg", "data": img_b64}},
                {"type": "text", "text": ANALYSIS_PROMPT},
            ],
        }],
    )
    raw = message.content[0].text.strip()

    start = raw.find("{")
    end = raw.rfind("}") + 1
    if start == -1:
        print("Erreur : pas de JSON dans la réponse")
        print(raw[:500])
        sys.exit(1)

    return json.loads(raw[start:end])


def apply_fill(shape, color_str):
    color = hex_to_rgb(color_str)
    if color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
    else:
        shape.fill.background()


def apply_line(shape, color_str, width_pt):
    color = hex_to_rgb(color_str)
    if color and width_pt and width_pt > 0:
        from pptx.util import Pt
        shape.line.color.rgb = color
        shape.line.width = Pt(width_pt)
    else:
        shape.line.fill.background()


def set_text(shape, elem):
    from pptx.util import Pt
    from pptx.enum.text import PP_ALIGN

    text = str(elem.get("text", "")).strip()
    if not text:
        return

    tf = shape.text_frame
    tf.word_wrap = True

    align_map = {
        "left": PP_ALIGN.LEFT,
        "center": PP_ALIGN.CENTER,
        "right": PP_ALIGN.RIGHT,
        "justify": PP_ALIGN.JUSTIFY,
    }

    # Vider les paragraphes existants
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align_map.get(str(elem.get("text_align", "left")).lower(), PP_ALIGN.LEFT)

    run = p.runs[0] if p.runs else p.add_run()
    run.font.bold = bool(elem.get("font_bold", False))
    run.font.italic = bool(elem.get("font_italic", False))
    run.font.size = Pt(float(elem.get("font_size_pt", 18)))

    font_color = hex_to_rgb(elem.get("font_color"))
    if font_color:
        run.font.color.rgb = font_color

    family = elem.get("font_family")
    if family and str(family).lower() not in ("null", "none", ""):
        run.font.name = str(family)


def build_pptx(analysis: dict, output_path: str):
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()

    aspect = str(analysis.get("slide", {}).get("aspect_ratio", "16:9"))
    if aspect == "4:3":
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
    else:  # 16:9
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)

    W = prs.slide_width
    H = prs.slide_height

    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # Fond du slide
    bg_color = hex_to_rgb(analysis.get("slide", {}).get("background_color", "#FFFFFF"))
    if bg_color:
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = bg_color

    elements = sorted(analysis.get("elements", []), key=lambda e: e.get("z_order", 0))
    print(f"  Création de {len(elements)} éléments...")

    for i, elem in enumerate(elements):
        x = int(float(elem.get("x_pct", 0)) / 100 * W)
        y = int(float(elem.get("y_pct", 0)) / 100 * H)
        w = max(int(float(elem.get("w_pct", 10)) / 100 * W), 914)   # min 1px
        h = max(int(float(elem.get("h_pct", 5)) / 100 * H), 914)

        etype = str(elem.get("type", "shape")).lower()
        stype = str(elem.get("shape_type", "rectangle")).lower()

        if etype == "text":
            shape = slide.shapes.add_textbox(x, y, w, h)
            shape.fill.background()
        else:
            shape_id = SHAPE_IDS.get(stype, 1)
            shape = slide.shapes.add_shape(shape_id, x, y, w, h)
            apply_fill(shape, elem.get("fill_color"))
            apply_line(shape, elem.get("border_color"), elem.get("border_width_pt", 0))

        if etype in ("text", "shape_with_text") or elem.get("text"):
            set_text(shape, elem)

    prs.save(output_path)
    print(f"\n  Sauvegardé : {output_path}")


def open_file(path: str):
    import subprocess
    subprocess.Popen(["start", "", path], shell=True)


def main():
    if len(sys.argv) < 2:
        print(__doc__)
        sys.exit(1)

    input_path = sys.argv[1]
    save_analysis = "--save-analysis" in sys.argv

    if input_path.endswith(".json"):
        with open(input_path, encoding="utf-8") as f:
            analysis = json.load(f)
        output_path = sys.argv[2] if len(sys.argv) > 2 else input_path.replace(".analysis.json", ".pptx")
        print(f"Reconstruction depuis analyse existante : {input_path}")
    else:
        if not Path(input_path).exists():
            print(f"Erreur : fichier introuvable : {input_path}")
            sys.exit(1)
        output_path = sys.argv[2] if len(sys.argv) > 2 else str(Path(input_path).with_suffix(".pptx"))

        print(f"Image -> PPTX : {input_path}")
        analysis = analyze_image(input_path)

        if save_analysis:
            analysis_path = Path(output_path).with_suffix("").with_suffix(".analysis.json")
            with open(analysis_path, "w", encoding="utf-8") as f:
                json.dump(analysis, f, indent=2, ensure_ascii=False)
            print(f"  Analyse sauvegardée : {analysis_path}")

    build_pptx(analysis, output_path)
    open_file(output_path)


if __name__ == "__main__":
    main()
