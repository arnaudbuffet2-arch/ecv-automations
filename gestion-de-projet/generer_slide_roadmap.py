import openpyxl
import datetime
import calendar
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# ─── COULEURS CHARTE (Rothschild) ─────────────────────────────────────────────
C_NAVY  = RGBColor(0x1C, 0x35, 0x5E)
C_SLATE = RGBColor(0x77, 0x86, 0x9E)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_DARK  = RGBColor(0x1A, 0x1A, 0x1A)
C_RED   = RGBColor(0xBF, 0x00, 0x00)
C_GRAY1 = RGBColor(0xF4, 0xF4, 0xF4)
C_GRAY2 = RGBColor(0xEB, 0xEB, 0xEB)
C_BORD  = RGBColor(0xC8, 0xC8, 0xC8)
C_TBD   = RGBColor(0xC0, 0xC0, 0xC0)

# Palette catégories
CAT_PALETTE = [
    '2E8B57',   # vert profond
    '4F86F7',   # bleu moderne
    '8B5CF6',   # violet doux
    'F59E0B',   # orange ambré
    '77869E',   # gris-bleu (5e si besoin)
]

def luminance(hexc):
    """Luminance relative (0=noir, 1=blanc)."""
    h = hexc.lstrip('#')
    r,g,b = int(h[0:2],16)/255, int(h[2:4],16)/255, int(h[4:6],16)/255
    return 0.299*r + 0.587*g + 0.114*b

def txt_on_bar(hexc):
    """Couleur texte adaptée à la couleur de fond de la barre."""
    return C_DARK if luminance(brand(hexc)) > 0.50 else C_WHITE

def h2c(hexc):
    h = hexc.lstrip('#')
    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

# ─── CHARGEMENT DONNÉES ───────────────────────────────────────────────────────

wb = openpyxl.load_workbook(r'C:\Users\arnau\OneDrive\Documents\roadmap_a_remplir autre.xlsx')

ws1 = wb['1 - Programme']
rows1 = [r for r in ws1.iter_rows(min_row=2, values_only=True) if any(c is not None for c in r)]
prog = rows1[-1]
TITLE    = str(prog[0]) if prog[0] else 'XXXXX'
SUBTITLE = str(prog[1]) if prog[1] else ''
END_YEAR = int(prog[3]) if prog[3] else 2027
today_raw = prog[4]
TODAY = today_raw.date() if isinstance(today_raw, datetime.datetime) else datetime.date(2026, 6, 1)

ws2 = wb['2 - Catégories']
categories = {}   # nom → hex (dans l'ordre Excel)
for row in ws2.iter_rows(min_row=2, values_only=True):
    if row[0] and row[1]:
        categories[row[0].strip()] = row[1].strip().lstrip('#')

# Remplacer les couleurs catégories par la palette distincte
cat_keys = list(categories.keys())
for i, k in enumerate(cat_keys):
    categories[k] = CAT_PALETTE[i % len(CAT_PALETTE)]

ws3 = wb['3 - Projets']
projects = {}
for row in ws3.iter_rows(min_row=2, values_only=True):
    if row[0] is not None and row[1]:
        num = int(row[0])
        projects[num] = {
            'num': num, 'name': str(row[1]),
            'category': str(row[3]).strip() if row[3] else '',
            'statut': str(row[5]) if row[5] else 'En cours',
            'phases': [], 'jalons': []
        }

def parse_date(val):
    if val is None or str(val).strip() == 'TBD': return None
    if isinstance(val, datetime.datetime):
        try:
            md = calendar.monthrange(val.year, val.month)[1]
            return datetime.date(val.year, val.month, min(val.day, md))
        except: return None
    if isinstance(val, str):
        try:
            p = val.strip().split('/')
            if len(p) == 3:
                d,m,y = int(p[0]),int(p[1]),int(p[2])
                md = calendar.monthrange(y,m)[1]
                return datetime.date(y,m,min(d,md))
        except: pass
    return None

ws4 = wb['4 - Phases des projets']
for row in ws4.iter_rows(min_row=2, values_only=True):
    if row[0] is not None and row[1]:
        num = int(row[0])
        if num in projects:
            s = parse_date(row[2])
            e = parse_date(row[3])
            if s is None:
                continue  # pas de date de début → on ignore
            is_tbd = (e is None)
            # Date de fin incohérente (fin < début) → TBD
            if e and e < s:
                e = None
                is_tbd = True
            projects[num]['phases'].append({
                'name': str(row[1]).strip(),
                'start': s, 'end': e, 'tbd': is_tbd
            })

ws5 = wb['5 - Jalons des projets']
for row in ws5.iter_rows(min_row=2, values_only=True):
    if row[0] is not None and row[1]:
        num = int(row[0])
        if num in projects:
            d = parse_date(row[2])
            if d: projects[num]['jalons'].append({'name': str(row[1]), 'date': d})

ws6 = wb['6 - Jalons globaux']
jalons_globaux = []
for row in ws6.iter_rows(min_row=2, values_only=True):
    if row[0] and row[1]:
        d = parse_date(row[1])
        if d: jalons_globaux.append({'name': str(row[0]), 'date': d})

# Correction : Lancement programme = Oct 2025 (saisi Oct 2026 par erreur)
for jg in jalons_globaux:
    if 'lancement' in jg['name'].lower():
        jg['date'] = datetime.date(2025, 10, 1)

# ─── TIMELINE ─────────────────────────────────────────────────────────────────

all_dates = [jg['date'] for jg in jalons_globaux]
for pr in projects.values():
    for ph in pr['phases']:
        if ph['start']: all_dates.append(ph['start'])

min_date = min(all_dates) if all_dates else datetime.date(2025, 10, 1)
TL_START = datetime.date(2026, 1, 1)
TL_END   = datetime.date(END_YEAR, 12, 31)

def gen_quarters(start, end):
    ql = ['Janv. – Mars','Avr. – Juin','Juil. – Sept.','Oct. – Déc.']
    qtrs, d = [], start
    while d <= end:
        q = (d.month - 1) // 3 + 1
        em = q * 3
        ed = calendar.monthrange(d.year, em)[1]
        qe = datetime.date(d.year, em, ed)
        qtrs.append({'year': d.year, 'q': q, 'label': f'T{q}',
                     'sublabel': ql[q-1], 'start': d, 'end': qe})
        nm = em + 1
        d = datetime.date(d.year + (1 if nm > 12 else 0), (nm-1)%12+1, 1)
    return qtrs

quarters = gen_quarters(TL_START, TL_END)
n_qtrs   = len(quarters)

# ─── LAYOUT ───────────────────────────────────────────────────────────────────

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)
HDR_H   = Inches(0.72)
NUM_X   = Inches(0.04);  NUM_W  = Inches(0.27)
NAME_X  = Inches(0.33);  NAME_W = Inches(2.02)
TL_X    = Inches(2.37);  TL_W   = Inches(10.78)
YEAR_H  = Inches(0.21)
QTR_H   = Inches(0.31)
ROW_H   = Inches(0.54)
HEADER_Y = HDR_H
QTR_Y    = HEADER_Y + YEAR_H
ROW_Y0   = QTR_Y + QTR_H
qtr_w    = TL_W / n_qtrs
BAR_H_1LANE = Inches(0.28)   # hauteur barre pour 1 lane (centrée dans la ligne)

# ─── POSITION PAR TRIMESTRE (alignement exact sur les colonnes) ───────────────

def dx(d):
    """Convertit une date en position X alignée sur les colonnes trimestre."""
    if d is None: return None
    if d <= TL_START: return TL_X
    if d >= TL_END:   return TL_X + TL_W
    for qi, q in enumerate(quarters):
        if q['start'] <= d <= q['end']:
            q_days = (q['end'] - q['start']).days + 1
            frac   = (d - q['start']).days / q_days
            return TL_X + (qi + frac) * qtr_w
    return TL_X + TL_W

# ─── HELPERS DESSIN ───────────────────────────────────────────────────────────

def rect(slide, x, y, w, h, fill=None, line=None, lw=Pt(0.5)):
    s = slide.shapes.add_shape(1, x, y, max(w, Emu(500)), max(h, Emu(500)))
    if fill: s.fill.solid(); s.fill.fore_color.rgb = fill
    else:    s.fill.background()
    if line: s.line.color.rgb = line; s.line.width = lw
    else:    s.line.fill.background()
    return s

def txt(slide, text, x, y, w, h, size=Pt(7), bold=False, color=None,
        align=PP_ALIGN.LEFT, wrap=True, italic=False):
    color = color or C_DARK
    tb = slide.shapes.add_textbox(x, y, max(w, Emu(500)), max(h, Emu(500)))
    tf = tb.text_frame; tf.word_wrap = wrap
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run(); run.text = text
    run.font.size = size; run.font.bold = bold
    run.font.color.rgb = color; run.font.italic = italic
    return tb

def oval(slide, x, y, w, h, fill, line=None):
    s = slide.shapes.add_shape(9, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line: s.line.color.rgb = line; s.line.width = Pt(0.75)
    else:    s.line.fill.background()
    return s

def diamond(slide, cx, cy, size, fill, line=None):
    s = slide.shapes.add_shape(1, cx - size/2, cy - size/2, size, size)
    s.rotation = 45.0
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line: s.line.color.rgb = line; s.line.width = Pt(0.75)
    else:    s.line.fill.background()
    return s

TBD_DAYS = 55   # durée affichée pour barres TBD

def assign_lanes(phases):
    """Affecte chaque phase à un couloir vertical pour éviter les superpositions."""
    n = len(phases)
    if n == 0: return [], 1
    assignments = [0] * n

    def eff_end(ph):
        s = ph['start']
        if not s: return datetime.date.max
        if ph['tbd'] or not ph['end']:
            return s + datetime.timedelta(days=TBD_DAYS)
        e = ph['end']
        return e if e > s else s + datetime.timedelta(days=1)

    order = sorted(range(n), key=lambda i: phases[i]['start'] or datetime.date.max)
    lane_ends = []
    for idx in order:
        ph = phases[idx]
        if not ph['start']: continue
        placed = False
        for li, le in enumerate(lane_ends):
            if ph['start'] >= le:
                lane_ends[li] = eff_end(ph)
                assignments[idx] = li
                placed = True; break
        if not placed:
            assignments[idx] = len(lane_ends)
            lane_ends.append(eff_end(ph))

    max_lane = max(assignments) if any(ph['start'] for ph in phases) else 0
    return assignments, max_lane + 1

CHARS_PER_INCH = 12.0   # approximation Pt(5), police étroite

def draw_bar_with_label(slide, label, x1, x2, bar_y, bar_h, fill_c, is_tbd):
    """Dessine la barre + texte correctement positionné."""
    bw = x2 - x1
    if bw < Emu(5000): bw = Emu(20000)

    b = rect(slide, x1, bar_y, bw, bar_h, fill=fill_c)
    if is_tbd:
        # Même couleur de fond, contour pointillé pour signaler "date de fin inconnue"
        b.line.color.rgb = fill_c
        b.line.width = Pt(1.5)
        ln = b.line._ln
        pd = etree.SubElement(ln, qn('a:prstDash'))
        pd.set('val', 'lgDash')

    # Largeur approximative du texte en pouces
    text_w_in = len(label) / CHARS_PER_INCH
    text_w = Inches(text_w_in)
    bw_in  = bw / 914400  # EMU → pouces

    # Luminance de la barre → texte blanc ou foncé selon le fond
    r_,g_,b_ = fill_c[0], fill_c[1], fill_c[2]
    lum = (0.299*r_ + 0.587*g_ + 0.114*b_) / 255
    bar_text_color = C_DARK if lum > 0.55 else C_WHITE

    if bw_in >= text_w_in * 1.15:
        # Texte à l'intérieur
        txt(slide, label, x1 + Inches(0.05), bar_y, bw - Inches(0.06), bar_h,
            size=Pt(5), color=bar_text_color, wrap=False, align=PP_ALIGN.LEFT)
    else:
        # Texte à l'extérieur, à droite si la place le permet
        out_x = x1 + bw + Inches(0.04)
        avail = (TL_X + TL_W) - out_x
        if avail > text_w * 0.6:
            txt(slide, label, out_x, bar_y, min(text_w + Inches(0.1), avail), bar_h,
                size=Pt(5), color=C_DARK, wrap=False, align=PP_ALIGN.LEFT)
        else:
            # À gauche si droite impossible
            out_x2 = x1 - text_w - Inches(0.04)
            if out_x2 >= TL_X:
                txt(slide, label, out_x2, bar_y, text_w, bar_h,
                    size=Pt(5), color=C_DARK, wrap=False, align=PP_ALIGN.RIGHT)
            # Sinon on tronque et on met dans la barre
            elif bw_in > Inches(0.4) / 914400:
                max_chars = max(3, int(bw_in * CHARS_PER_INCH) - 2)
                txt(slide, label[:max_chars] + '…', x1 + Inches(0.03), bar_y,
                    bw - Inches(0.04), bar_h,
                    size=Pt(5), color=C_WHITE, wrap=False, align=PP_ALIGN.LEFT)

# ─── CONSTRUCTION SLIDE ───────────────────────────────────────────────────────

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H
slide = prs.slides.add_slide(prs.slide_layouts[6])

# ── En-tête ───────────────────────────────────────────────────────────────────
rect(slide, Inches(0), Inches(0), SLIDE_W, HDR_H, fill=C_NAVY)
txt(slide, TITLE, Inches(0.12), Inches(0.03), Inches(7.1), Inches(0.39),
    size=Pt(17), bold=True, color=C_WHITE, align=PP_ALIGN.LEFT)
txt(slide, SUBTITLE, Inches(0.12), Inches(0.40), Inches(7.1), Inches(0.27),
    size=Pt(9.5), color=C_SLATE, align=PP_ALIGN.LEFT)

# ── Légende : Types de projets ────────────────────────────────────────────────
LX, LY, LW = Inches(7.5), Inches(0.05), Inches(2.85)
valid_cats = [(k, v) for k, v in categories.items() if k and v]
LH = Inches(0.17 + 0.09 * len(valid_cats) + 0.04)
rect(slide, LX, LY, LW, LH, fill=RGBColor(0x2A,0x45,0x74))
txt(slide, 'TYPES DE PROJETS', LX+Inches(0.06), LY+Inches(0.02),
    LW-Inches(0.1), Inches(0.14), size=Pt(6), bold=True, color=C_WHITE)
for i,(name, hexc) in enumerate(valid_cats):
    ry = LY + Inches(0.17) + i * Inches(0.09)
    rect(slide, LX+Inches(0.07), ry+Inches(0.01), Inches(0.13), Inches(0.068), fill=h2c(hexc))
    txt(slide, name, LX+Inches(0.24), ry, LW-Inches(0.28), Inches(0.088), size=Pt(5.5), color=C_WHITE)

# ── Légende : Conventions ────────────────────────────────────────────────────
SX, SY, SW = LX+LW+Inches(0.07), LY, Inches(2.30)
SH = Inches(0.17 + 0.09 * 2 + 0.04)
rect(slide, SX, SY, SW, SH, fill=RGBColor(0x2A,0x45,0x74))
txt(slide, 'CONVENTIONS', SX+Inches(0.06), SY+Inches(0.02),
    SW-Inches(0.1), Inches(0.14), size=Pt(6), bold=True, color=C_WHITE)
convs = [
    ('Dates précises',            h2c('2E8B57'), False),
    ('Date de fin non précisée',  h2c('2E8B57'), True),   # même couleur + bordure pointillée
]
for i,(label, color, tbd_sample) in enumerate(convs):
    ry = SY + Inches(0.17) + i * Inches(0.09)
    s = rect(slide, SX+Inches(0.07), ry+Inches(0.01), Inches(0.13), Inches(0.068), fill=color)
    if tbd_sample:
        s.line.color.rgb = color
        s.line.width = Pt(1.5)
        ln = s.line._ln
        pd = etree.SubElement(ln, qn('a:prstDash'))
        pd.set('val', 'lgDash')
    txt(slide, label, SX+Inches(0.24), ry, SW-Inches(0.28), Inches(0.088), size=Pt(5.5), color=C_WHITE)

# ── En-têtes colonnes (#, PROJETS) ────────────────────────────────────────────
rect(slide, Inches(0), HEADER_Y, TL_X, YEAR_H+QTR_H, fill=C_NAVY)
txt(slide, '#', NUM_X, HEADER_Y, NUM_W, YEAR_H+QTR_H,
    size=Pt(8), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
txt(slide, 'PROJETS', NAME_X, HEADER_Y, NAME_W, YEAR_H+QTR_H,
    size=Pt(8), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

# ── En-têtes années ───────────────────────────────────────────────────────────
from itertools import groupby
for year, grp in groupby(quarters, key=lambda q: q['year']):
    grp = list(grp)
    xs = dx(grp[0]['start'])
    xe = dx(grp[-1]['end']) + qtr_w * 0   # fin = début du trimestre suivant
    # Plus précis : fin = début trimestre suivant OU TL_X+TL_W
    xs = max(xs, TL_X)
    xe = dx(grp[-1]['end']) + Emu(1)  # +1 EMU pour inclure la fin exacte
    xe = min(TL_X + TL_W, dx(grp[-1]['end']) + qtr_w * (1/max((grp[-1]['end']-grp[-1]['start']).days,1)))
    # Plus simple : utiliser les indices
    qi_start = quarters.index(grp[0])
    qi_end   = quarters.index(grp[-1])
    xs = TL_X + qi_start * qtr_w
    xe = TL_X + (qi_end + 1) * qtr_w
    w = xe - xs
    rect(slide, xs, HEADER_Y, w, YEAR_H, fill=C_NAVY, line=C_WHITE, lw=Pt(0.5))
    txt(slide, str(year), xs, HEADER_Y, w, YEAR_H,
        size=Pt(8.5), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

# ── En-têtes trimestres ───────────────────────────────────────────────────────
for qi, q in enumerate(quarters):
    x = TL_X + qi * qtr_w
    rect(slide, x, QTR_Y, qtr_w, QTR_H,
         fill=RGBColor(0x3A,0x5A,0x80), line=C_WHITE, lw=Pt(0.4))
    txt(slide, q['label'], x, QTR_Y, qtr_w, Inches(0.15),
        size=Pt(6.5), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    txt(slide, q['sublabel'], x, QTR_Y+Inches(0.15), qtr_w, Inches(0.15),
        size=Pt(5.3), color=C_WHITE, align=PP_ALIGN.CENTER)

# ── Lignes de séparation verticales (trim) ────────────────────────────────────
sorted_proj = sorted(projects.values(), key=lambda p: p['num'])
n_rows = len(sorted_proj)
rows_area_h = n_rows * ROW_H
for qi in range(1, n_qtrs):
    vx = TL_X + qi * qtr_w
    rect(slide, vx - Inches(0.002), ROW_Y0, Inches(0.004), rows_area_h,
         fill=RGBColor(0xDD,0xDD,0xDD))

# ── Lignes projets ────────────────────────────────────────────────────────────
for i, proj in enumerate(sorted_proj):
    ry = ROW_Y0 + i * ROW_H
    bg = C_GRAY1 if i % 2 == 0 else C_GRAY2
    rect(slide, Inches(0), ry, SLIDE_W, ROW_H, fill=bg, line=C_BORD, lw=Pt(0.25))

    # Cellule #
    rect(slide, NUM_X, ry, NUM_W, ROW_H, fill=C_NAVY, line=C_BORD, lw=Pt(0.25))
    txt(slide, str(proj['num']), NUM_X, ry, NUM_W, ROW_H,
        size=Pt(7.5), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    # Cellule nom
    rect(slide, NAME_X, ry, NAME_W+Inches(0.02), ROW_H, line=C_BORD, lw=Pt(0.25))
    txt(slide, proj['name'], NAME_X+Inches(0.04), ry+Inches(0.03),
        NAME_W-Inches(0.04), ROW_H-Inches(0.06),
        size=Pt(6), color=C_DARK, wrap=True)

    cat_hex = categories.get(proj['category'], '77869E')
    bar_color = h2c(cat_hex)

    phases = proj['phases']
    assignments, n_lanes = assign_lanes(phases)

    # Calcul de la géométrie des lanes
    MARGIN   = Inches(0.055)
    LANE_GAP = Inches(0.02)
    avail_h  = ROW_H - 2 * MARGIN

    if n_lanes == 1:
        lane_h = min(BAR_H_1LANE, avail_h)
        bar_y_base = ry + (ROW_H - lane_h) / 2   # centré
    else:
        lane_h = (avail_h - LANE_GAP * (n_lanes - 1)) / n_lanes
        bar_y_base = ry + MARGIN

    for pi, phase in enumerate(phases):
        if not phase['start']: continue
        lane  = assignments[pi]
        bar_y = bar_y_base if n_lanes == 1 else bar_y_base + lane * (lane_h + LANE_GAP)
        bar_h = lane_h

        s_date = phase['start']
        e_date = phase['end']
        is_tbd = phase['tbd']

        # Phase ponctuelle (même date début = fin) → losange
        if not is_tbd and e_date and e_date == s_date:
            mx = dx(s_date)
            if mx:
                ds = min(bar_h * 0.85, Inches(0.14))
                diamond(slide, mx, bar_y + bar_h/2, ds, fill=bar_color)
                # Étiquette au-dessus
                lw = Inches(0.9)
                txt(slide, phase['name'], mx - lw/2, bar_y - Inches(0.15),
                    lw, Inches(0.14), size=Pt(4.5), color=C_DARK,
                    align=PP_ALIGN.CENTER, wrap=False)
            continue

        x1 = dx(s_date)
        if x1 is None: continue

        if is_tbd:
            x2 = dx(s_date + datetime.timedelta(days=TBD_DAYS))
        else:
            x2 = dx(e_date)
        if x2 is None: x2 = x1 + Inches(0.5)

        x1 = max(x1, TL_X)
        x2 = min(x2, TL_X + TL_W)

        # TBD = même couleur catégorie, mais bord pointillé (géré dans draw_bar_with_label)
        draw_bar_with_label(slide, phase['name'], x1, x2, bar_y, bar_h, bar_color, is_tbd)

    # Jalons du projet (losanges depuis onglet 5)
    ref_y = bar_y_base + (lane_h / 2)
    for jal in proj['jalons']:
        jx = dx(jal['date'])
        if not jx: continue
        ds = Inches(0.13)
        diamond(slide, jx, ref_y, ds, fill=C_NAVY)
        lw = Inches(0.8)
        txt(slide, jal['name'], jx - lw/2, ref_y - Inches(0.17),
            lw, Inches(0.15), size=Pt(4.5), color=C_DARK,
            align=PP_ALIGN.CENTER, wrap=False)

# ── Ligne TODAY ───────────────────────────────────────────────────────────────
rows_end_y = ROW_Y0 + n_rows * ROW_H
today_x = dx(TODAY)
if today_x and TL_X <= today_x <= TL_X + TL_W:
    rect(slide, today_x - Inches(0.007), HEADER_Y,
         Inches(0.014), rows_end_y - HEADER_Y, fill=C_RED)
    lw = Inches(0.55)
    rect(slide, today_x - lw/2, rows_end_y, lw, Inches(0.19), fill=C_RED)
    txt(slide, 'TODAY', today_x - lw/2, rows_end_y, lw, Inches(0.19),
        size=Pt(6.5), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

# ── Jalons clés globaux (barre du bas) ────────────────────────────────────────
BOT_Y = rows_end_y + Inches(0.23)
BOT_H = Inches(0.88)
rect(slide, Inches(0), BOT_Y, SLIDE_W, BOT_H, fill=C_NAVY)
txt(slide, 'JALONS CLÉS', NUM_X, BOT_Y, TL_X-Inches(0.08), BOT_H,
    size=Pt(7.5), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

line_y = BOT_Y + BOT_H * 0.5
rect(slide, TL_X, line_y - Inches(0.007), TL_W, Inches(0.014), fill=C_SLATE)

for jg in jalons_globaux:
    jx = dx(jg['date'])
    if not jx: continue
    ds = Inches(0.14)
    diamond(slide, jx, line_y, ds, fill=C_WHITE, line=C_NAVY)
    lw = Inches(0.95)
    txt(slide, jg['name'], jx - lw/2, line_y - Inches(0.33),
        lw, Inches(0.18), size=Pt(5.5), color=C_WHITE,
        align=PP_ALIGN.CENTER, wrap=False)
    txt(slide, jg['date'].strftime('%b. %Y'), jx - lw/2, line_y+Inches(0.10),
        lw, Inches(0.16), size=Pt(5), color=C_SLATE,
        align=PP_ALIGN.CENTER, wrap=False)

# ─── SAUVEGARDE ───────────────────────────────────────────────────────────────
out = r'c:\cerveau 2 obsidian vault\ceveau 2 vault\data\roadmap_output_v5.pptx'
prs.save(out)
print('OK :', out)
