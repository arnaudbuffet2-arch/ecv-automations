#!/usr/bin/env python3
"""
Slide hybride : PNG haute-résolution en fond + zones de texte éditables par-dessus.
Chaque bloc de texte est cliquable/modifiable dans PowerPoint.
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

NAVY  = RGBColor(0x1B, 0x2A, 0x4A)
GOLD  = RGBColor(0xC4, 0x8C, 0x28)
MGREY = RGBColor(0x8A, 0x9B, 0xB5)
LGREY = RGBColor(0xCC, 0xD2, 0xDA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

PNG = Path(r"C:\Users\arnau\OneDrive\Bureau\slide_ath_render.png")
OUT = Path(r"C:\Users\arnau\OneDrive\Bureau\slide_ath_editeur.pptx")

FONT_TITLE = "Calibri Light"
FONT_BODY  = "Calibri"


def T(sl, text, x, y, w, h, size, bold=False, italic=False,
      color=NAVY, align=PP_ALIGN.LEFT, font=None):
    font = font or FONT_BODY
    tb = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for line in text.split('\n'):
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color
        r.font.name = font
    # Transparent background, no border
    tb.fill.background()
    tb.line.fill.background()
    return tb


prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.50)
sl = prs.slides.add_slide(prs.slide_layouts[6])

# ── Background image (full slide)
sl.shapes.add_picture(str(PNG), Inches(0), Inches(0), Inches(13.33), Inches(7.50))

# ── TITLE
T(sl, "Processus de sélection de la solution cible",
  0.25, 0.10, 12.3, 0.66, 26, bold=True, font=FONT_TITLE)

# ── SUBTITLE
T(sl, "Démarche d'évaluation des éditeurs et aide à la décision",
  0.25, 0.80, 11.0, 0.38, 12, italic=True)

# ── PROCESS COLUMNS
CX = [0.25, 3.42, 6.59, 9.76]
CW = 2.95

COLS = [
    ("CADRAGE", [
        "Finalisation du besoin\net du contenu du RFP",
        "Identification de 10 éditeurs\ncibles",
        "Envoi du RFP et calendrier\nassocié",
    ]),
    ("ÉVALUATION", [
        "Définition des critères d'évaluation\n(fonctionnel, IT, coût, références...)",
        "Pondération par catégorie",
        "Formalisation de la matrice\nd'évaluation",
        "Gestion des questions/réponses\néditeurs",
    ]),
    ("SÉLECTION", [
        "Scoring des réponses\n(fonctionnel, financier, planning...)",
        "Shortlist de 2–3 éditeurs",
        "Organisation des soutenances /\ndémos / POC",
        "Consolidation des évaluations",
    ]),
    ("DÉCISION & DÉPLOIEMENT", [
        "Recommandation finale",
        "Construction du macro–planning\n(POC, migration, budget)",
        "Validation en comité\nde décision",
    ]),
]

for i, (title, bullets) in enumerate(COLS):
    cx = CX[i]
    T(sl, title, cx, 1.98, CW, 0.34, 8.5, bold=True, align=PP_ALIGN.CENTER)
    for k, b in enumerate(bullets):
        T(sl, "▪  " + b, cx + 0.07, 2.47 + k * 0.37, CW - 0.12, 0.35, 7.5)

# ── RÔLE DU CABINET
T(sl, "RÔLE DU CABINET\nDE CONSEIL", 0.72, 4.14, 1.55, 0.55, 7.5, bold=True)

ROLES = [
    ("Structurer la démarche\net cadrer le besoin",                    2.30, 2.75),
    ("Construire le cadre\nd'évaluation",                               5.10, 1.90),
    ("Soutien opérationnel\nchef de projet IT +\nchef de projet Métier", 7.05, 3.00),
    ("Consolider les analyses\net formuler la recommandation",          10.10, 2.95),
]
for txt, rx, rw in ROLES:
    T(sl, txt, rx, 4.14, rw, 0.56, 7.5)

# ── CRITÈRES D'ÉVALUATION
T(sl, "CRITÈRES D'ÉVALUATION", 0.25, 4.76, 3.5, 0.22, 7.5, bold=True)

HX = [3.20, 4.10, 5.00]
for j, h in enumerate(["PRESTATAIRE A", "PRESTATAIRE B", "PRESTATAIRE C"]):
    T(sl, h, HX[j], 5.02, 0.88, 0.21, 5.5, bold=True, align=PP_ALIGN.CENTER)

CRIT = [
    "Couverture fonctionnelle",
    "Ergonomie & expérience utilisateur",
    "Architecture & intégration SI",
    "Références & expérience marché",
    "Coût total de possession",
    "Planning & capacité de déploiement",
    "Accompagnement au changement",
    "Sécurité & conformité réglementaire",
    "Roadmap produit & innovation",
]
RH = 0.194
for k, crit in enumerate(CRIT):
    T(sl, crit, 0.48, 5.25 + k * RH, 2.68, RH, 7.0)

# ── POSITIONNEMENT
T(sl, "POSITIONNEMENT DES SOLUTIONS ÉVALUÉES",
  6.30, 4.76, 4.80, 0.22, 7.5, bold=True)

# ── RECOMMANDATION BOX
T(sl, "RECOMMANDATION",
  10.30, 5.03, 2.70, 0.24, 8.0, bold=True, align=PP_ALIGN.CENTER)
T(sl, "SOLUTION\nRECOMMANDÉE",
  10.30, 5.60, 2.70, 0.46, 9.5, bold=True, align=PP_ALIGN.CENTER)
T(sl, "PRESTATAIRE A",
  10.30, 6.09, 2.70, 0.30, 13.0, bold=True, align=PP_ALIGN.CENTER, color=GOLD)
T(sl,
  "Le Prestataire A est recommandé, étant le seul à couvrir l'ensemble "
  "des besoins prioritaires identifiés lors du cadrage, avec une adéquation "
  "optimale aux enjeux fonctionnels, techniques et organisationnels.",
  10.28, 6.44, 2.74, 0.76, 7.0)

# ── FOOTER
T(sl,
  "Soutien opérationnel chef de projet IT  +  Chef de projet Métier  +  Métier",
  0.25, 7.33, 12.83, 0.18, 8.0, italic=True, align=PP_ALIGN.CENTER)

prs.save(str(OUT))
print(f"Sauvegardé : {OUT}")
print("Ouvrir dans PowerPoint → double-clic sur n'importe quel texte pour l'éditer.")
