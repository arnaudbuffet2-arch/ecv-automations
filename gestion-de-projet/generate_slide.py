#!/usr/bin/env python3
"""
generate_slide.py — Converts a slide image (PNG/JPG) to an editable PowerPoint file.
Uses OpenCV for layout analysis and pytesseract for OCR (optional).
"""

import argparse
import os
import subprocess
from datetime import datetime
from pathlib import Path
from typing import Optional

import cv2
import numpy as np
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt

try:
    import pytesseract
    # Auto-detect Tesseract on Windows if not in PATH
    _WIN_PATHS = [
        r"C:\Program Files\Tesseract-OCR\tesseract.exe",
        r"C:\Program Files (x86)\Tesseract-OCR\tesseract.exe",
    ]
    if os.name == "nt":
        for _p in _WIN_PATHS:
            if os.path.exists(_p):
                pytesseract.pytesseract.tesseract_cmd = _p
                break
    pytesseract.get_tesseract_version()
    TESSERACT_AVAILABLE = True
except Exception:
    TESSERACT_AVAILABLE = False

# Slide dimensions in EMU (English Metric Units): 1 inch = 914400 EMU
SLIDE_16_9 = (12192000, 6858000)
SLIDE_4_3 = (9144000, 6858000)
EMU_PER_PT = 12700


# ---------------------------------------------------------------------------
# Image loading
# ---------------------------------------------------------------------------

def load_image(path: str) -> tuple:
    """Load image from disk, return (cv2 BGR ndarray, PIL Image RGB)."""
    if not os.path.exists(path):
        raise FileNotFoundError(f"Image not found: {path}")
    img_cv = cv2.imread(path)
    if img_cv is None:
        raise ValueError(f"Cannot read image (unsupported format?): {path}")
    img_pil = Image.open(path).convert("RGB")
    return img_cv, img_pil


# ---------------------------------------------------------------------------
# Layout analysis
# ---------------------------------------------------------------------------

def detect_slide_ratio(img: np.ndarray) -> tuple:
    """Return (ratio_float, ratio_label, slide_dims_emu)."""
    h, w = img.shape[:2]
    ratio = w / h
    if abs(ratio - 16 / 9) < 0.15:
        return ratio, "16:9", SLIDE_16_9
    if abs(ratio - 4 / 3) < 0.10:
        return ratio, "4:3", SLIDE_4_3
    # Custom ratio: preserve proportions based on 16:9 height
    height_emu = SLIDE_16_9[1]
    width_emu = int(height_emu * ratio)
    return ratio, f"custom ({w}×{h})", (width_emu, height_emu)


def extract_dominant_colors(img: np.ndarray, k: int = 6) -> list:
    """K-means dominant colors. Returns list of (R, G, B) tuples."""
    pixels = img.reshape(-1, 3).astype(np.float32)
    criteria = (cv2.TERM_CRITERIA_EPS + cv2.TERM_CRITERIA_MAX_ITER, 20, 1.0)
    _, labels, centers = cv2.kmeans(pixels, k, None, criteria, 5, cv2.KMEANS_RANDOM_CENTERS)
    # Sort by frequency (most common first)
    counts = np.bincount(labels.flatten())
    order = np.argsort(-counts)
    centers = centers.astype(int)[order]
    return [(int(c[2]), int(c[1]), int(c[0])) for c in centers]  # BGR → RGB


def detect_visual_blocks(img: np.ndarray) -> list:
    """
    Detect significant visual regions via contour analysis.
    Returns list of (x, y, w, h) in image pixels, sorted top-to-bottom.
    """
    gray = cv2.cvtColor(img, cv2.COLOR_BGR2GRAY)
    blurred = cv2.GaussianBlur(gray, (5, 5), 0)
    edges = cv2.Canny(blurred, 30, 100)
    kernel = np.ones((7, 7), np.uint8)
    dilated = cv2.dilate(edges, kernel, iterations=2)

    contours, _ = cv2.findContours(dilated, cv2.RETR_EXTERNAL, cv2.CHAIN_APPROX_SIMPLE)
    img_h, img_w = img.shape[:2]
    min_area = img_w * img_h * 0.004

    blocks = []
    for cnt in contours:
        x, y, w, h = cv2.boundingRect(cnt)
        if w * h > min_area and w > 15 and h > 8:
            if w > img_w * 0.95 and h > img_h * 0.95:
                continue  # Skip full-image contour
            blocks.append((x, y, w, h))

    blocks = _merge_blocks(blocks)
    return sorted(blocks, key=lambda b: (b[1], b[0]))


def _merge_blocks(blocks: list) -> list:
    """Iteratively merge overlapping bounding boxes."""
    if not blocks:
        return blocks
    merged = True
    while merged:
        merged = False
        used = set()
        result = []
        for i, b1 in enumerate(blocks):
            if i in used:
                continue
            x1, y1, w1, h1 = b1
            r1, bot1 = x1 + w1, y1 + h1
            for j in range(i + 1, len(blocks)):
                if j in used:
                    continue
                x2, y2, w2, h2 = blocks[j]
                r2, bot2 = x2 + w2, y2 + h2
                if max(x1, x2) < min(r1, r2) and max(y1, y2) < min(bot1, bot2):
                    x1, y1 = min(x1, x2), min(y1, y2)
                    r1, bot1 = max(r1, r2), max(bot1, bot2)
                    w1, h1 = r1 - x1, bot1 - y1
                    used.add(j)
                    merged = True
            result.append((x1, y1, w1, h1))
            used.add(i)
        blocks = result
    return blocks


def detect_lines(img: np.ndarray) -> tuple:
    """
    Detect horizontal and vertical lines via Hough transform.
    Returns (h_lines, v_lines) as lists of (x1, y1, x2, y2).
    """
    gray = cv2.cvtColor(img, cv2.COLOR_BGR2GRAY)
    edges = cv2.Canny(gray, 50, 150, apertureSize=3)
    lines = cv2.HoughLinesP(edges, 1, np.pi / 180, threshold=80,
                             minLineLength=img.shape[1] // 10, maxLineGap=15)
    h_lines, v_lines = [], []
    if lines is not None:
        for seg in lines:
            x1, y1, x2, y2 = seg[0]
            angle = abs(np.degrees(np.arctan2(y2 - y1, x2 - x1)))
            if angle < 5 or angle > 175:
                h_lines.append((min(x1, x2), (y1 + y2) // 2, max(x1, x2), (y1 + y2) // 2))
            elif 85 < angle < 95:
                v_lines.append(((x1 + x2) // 2, min(y1, y2), (x1 + x2) // 2, max(y1, y2)))
    return h_lines, v_lines


def detect_text_zones(img: np.ndarray) -> list:
    """
    Detect text-like regions using adaptive threshold + morphology.
    Returns list of (x, y, w, h), sorted top-to-bottom.
    """
    gray = cv2.cvtColor(img, cv2.COLOR_BGR2GRAY)
    thresh = cv2.adaptiveThreshold(gray, 255, cv2.ADAPTIVE_THRESH_GAUSSIAN_C,
                                   cv2.THRESH_BINARY_INV, 15, 4)
    kernel = np.ones((3, 25), np.uint8)
    dilated = cv2.dilate(thresh, kernel, iterations=1)
    contours, _ = cv2.findContours(dilated, cv2.RETR_EXTERNAL, cv2.CHAIN_APPROX_SIMPLE)
    img_w = img.shape[1]
    zones = []
    for cnt in contours:
        x, y, w, h = cv2.boundingRect(cnt)
        if w > 30 and h > 6 and w < img_w * 0.98 and w / h > 1.5:
            zones.append((x, y, w, h))
    return sorted(zones, key=lambda z: (z[1], z[0]))


# ---------------------------------------------------------------------------
# OCR
# ---------------------------------------------------------------------------

def ocr_image(img_pil: Image.Image) -> list:
    """
    Run Tesseract OCR. Returns list of word dicts with text + bounding box.
    Returns empty list if Tesseract is unavailable.
    """
    if not TESSERACT_AVAILABLE:
        return []
    try:
        data = pytesseract.image_to_data(img_pil, output_type=pytesseract.Output.DICT,
                                          lang="fra+eng")
        results = []
        for i in range(len(data["text"])):
            text = data["text"][i].strip()
            conf = int(data["conf"][i])
            if text and conf > 35:
                results.append({
                    "text": text,
                    "x": data["left"][i], "y": data["top"][i],
                    "w": data["width"][i], "h": data["height"][i],
                    "conf": conf,
                })
        return results
    except Exception:
        return []


def group_ocr_into_lines(words: list, gap: int = 12) -> list:
    """Merge individual OCR words into text lines."""
    if not words:
        return []
    sorted_words = sorted(words, key=lambda w: (w["y"], w["x"]))
    lines, current = [], [sorted_words[0]]
    for word in sorted_words[1:]:
        if abs(word["y"] - current[-1]["y"]) <= gap:
            current.append(word)
        else:
            lines.append(current)
            current = [word]
    lines.append(current)

    merged = []
    for line in lines:
        line = sorted(line, key=lambda w: w["x"])
        x = line[0]["x"]
        y = min(w["y"] for w in line)
        right = max(w["x"] + w["w"] for w in line)
        bottom = max(w["y"] + w["h"] for w in line)
        merged.append({
            "text": " ".join(w["text"] for w in line),
            "x": x, "y": y, "w": right - x, "h": bottom - y,
        })
    return merged


# ---------------------------------------------------------------------------
# Coordinate conversion & color helpers
# ---------------------------------------------------------------------------

def coords_to_emu(x: int, y: int, w: int, h: int,
                  img_w: int, img_h: int, slide_dims: tuple) -> tuple:
    """Scale image pixel coords to PowerPoint EMU coords."""
    sw, sh = slide_dims
    return int(x * sw / img_w), int(y * sh / img_h), int(w * sw / img_w), int(h * sh / img_h)


def get_region_color(img: np.ndarray, x: int, y: int, w: int, h: int) -> tuple:
    """Average color of a region, returned as (R, G, B)."""
    roi = img[max(0, y):max(0, y) + max(1, h), max(0, x):max(0, x) + max(1, w)]
    if roi.size == 0:
        return (128, 128, 128)
    avg = roi.mean(axis=(0, 1))
    return (int(avg[2]), int(avg[1]), int(avg[0]))  # BGR → RGB


def _luminance(rgb: tuple) -> float:
    r, g, b = rgb
    return 0.299 * r + 0.587 * g + 0.114 * b


def _font_color(bg: tuple) -> tuple:
    return (255, 255, 255) if _luminance(bg) < 128 else (20, 20, 20)


def _estimate_font_pt(h_emu: int) -> int:
    raw = int(h_emu / EMU_PER_PT * 0.72)
    return max(8, min(raw, 72))


# ---------------------------------------------------------------------------
# PPTX construction
# ---------------------------------------------------------------------------

def build_pptx(img: np.ndarray, slide_dims: tuple, dominant_colors: list,
               blocks: list, h_lines: list, v_lines: list,
               ocr_lines: list) -> Presentation:
    """Assemble the PowerPoint slide from detected elements."""
    prs = Presentation()
    sw, sh = slide_dims
    prs.slide_width, prs.slide_height = sw, sh

    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    img_h, img_w = img.shape[:2]

    # --- Background ---
    bg_color = dominant_colors[0] if dominant_colors else (255, 255, 255)
    bg = slide.shapes.add_shape(1, 0, 0, sw, sh)  # 1 = rectangle
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(*bg_color)
    bg.line.width = 0

    # --- Visual blocks ---
    for (x, y, w, h) in blocks:
        ex, ey, ew, eh = coords_to_emu(x, y, w, h, img_w, img_h, slide_dims)
        if ew < 5000 or eh < 2000:
            continue
        color = get_region_color(img, x, y, w, h)
        if sum(abs(color[i] - bg_color[i]) for i in range(3)) < 18:
            continue  # Skip near-invisible blocks
        shape = slide.shapes.add_shape(1, ex, ey, ew, eh)
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(*color)
        shape.line.color.rgb = RGBColor(*color)

    # --- Lines ---
    for (x1, y1, x2, y2) in h_lines + v_lines:
        ex1, ey1, _, _ = coords_to_emu(x1, y1, 1, 1, img_w, img_h, slide_dims)
        ex2, ey2, _, _ = coords_to_emu(x2, y2, 1, 1, img_w, img_h, slide_dims)
        if abs(ex2 - ex1) < 500 and abs(ey2 - ey1) < 500:
            continue
        try:
            conn = slide.shapes.add_connector(1, ex1, ey1, ex2, ey2)
            lc = get_region_color(img, x1, y1, max(1, abs(x2 - x1)), max(1, abs(y2 - y1)))
            conn.line.color.rgb = RGBColor(*lc)
            conn.line.width = Pt(1.5)
        except Exception:
            pass

    # --- OCR text boxes ---
    for item in ocr_lines:
        x, y, w, h = item["x"], item["y"], item["w"], item["h"]
        pad = 5
        ex, ey, ew, eh = coords_to_emu(
            max(0, x - pad), max(0, y - pad), w + 2 * pad, h + 2 * pad,
            img_w, img_h, slide_dims
        )
        if ew < 20000 or eh < 8000:
            continue
        bg_at_text = get_region_color(img, x, y, w, h)
        fc = _font_color(bg_at_text)
        fsize = _estimate_font_pt(eh)
        txbox = slide.shapes.add_textbox(ex, ey, ew, eh)
        tf = txbox.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = item["text"]
        run.font.size = Pt(fsize)
        run.font.color.rgb = RGBColor(*fc)
        # Heuristic: bold if large text near top
        if ey < sh * 0.25 and fsize > 16:
            run.font.bold = True

    return prs


# ---------------------------------------------------------------------------
# Preview & debug
# ---------------------------------------------------------------------------

def export_preview(pptx_path: str, preview_path: str) -> bool:
    """Try to convert PPTX to PNG using LibreOffice. Returns True on success."""
    out_dir = str(Path(preview_path).parent)
    candidates = [
        "soffice", "libreoffice",
        r"C:\Program Files\LibreOffice\program\soffice.exe",
        r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
        "/usr/bin/soffice", "/usr/local/bin/soffice",
    ]
    for exe in candidates:
        try:
            r = subprocess.run(
                [exe, "--headless", "--convert-to", "png", "--outdir", out_dir, pptx_path],
                capture_output=True, timeout=30
            )
            if r.returncode == 0:
                generated = Path(out_dir) / (Path(pptx_path).stem + ".png")
                if generated.exists():
                    if str(generated) != preview_path:
                        generated.rename(preview_path)
                    return True
        except (FileNotFoundError, subprocess.TimeoutExpired, OSError):
            continue
    return False


def save_debug_image(img: np.ndarray, blocks: list, h_lines: list, v_lines: list,
                     ocr_lines: list, output_path: str) -> None:
    """Save annotated image showing all detected elements."""
    dbg = img.copy()
    for (x, y, w, h) in blocks:
        cv2.rectangle(dbg, (x, y), (x + w, y + h), (0, 200, 0), 2)
    for (x1, y1, x2, y2) in h_lines:
        cv2.line(dbg, (x1, y1), (x2, y2), (255, 80, 0), 2)
    for (x1, y1, x2, y2) in v_lines:
        cv2.line(dbg, (x1, y1), (x2, y2), (0, 80, 255), 2)
    for item in ocr_lines:
        x, y, w, h = item["x"], item["y"], item["w"], item["h"]
        cv2.rectangle(dbg, (x, y), (x + w, y + h), (0, 165, 255), 1)
    cv2.imwrite(output_path, dbg)


# ---------------------------------------------------------------------------
# Markdown report
# ---------------------------------------------------------------------------

def generate_report(input_path: str, img_w: int, img_h: int, ratio_name: str,
                    dominant_colors: list, blocks: list, text_zones: list,
                    ocr_lines: list, ocr_ok: bool, preview_ok: bool,
                    debug: bool = False) -> str:
    lines = [
        "# Rapport de reconstruction slide",
        "",
        f"**Source** : `{input_path}`",
        f"**Date** : {datetime.now().strftime('%Y-%m-%d %H:%M')}",
        "",
        "## Analyse",
        "",
        f"| Paramètre | Valeur |",
        f"|-----------|--------|",
        f"| Dimensions | {img_w} × {img_h} px |",
        f"| Ratio détecté | {ratio_name} |",
        f"| Blocs visuels | {len(blocks)} |",
        f"| Zones de texte | {len(text_zones)} |",
        f"| Lignes OCR | {len(ocr_lines)} |",
        "",
        "## Couleurs dominantes",
        "",
    ]
    for i, (r, g, b) in enumerate(dominant_colors[:5]):
        lines.append(f"- `rgb({r}, {g}, {b})`")
    lines += [
        "",
        "## Statut",
        "",
        f"- OCR Tesseract : {'✅ disponible' if ocr_ok else '❌ non disponible — mode dégradé (textes non extraits)'}",
        f"- Aperçu PNG : {'✅ généré' if preview_ok else '❌ non généré — LibreOffice requis (`soffice --headless`)'}",
        "",
        "## Limites connues",
        "",
        "- Polices exactes non détectées (approximées)",
        "- Effets visuels (ombres, dégradés) non reproduits",
        "- Icônes complexes approximées en formes simples",
        "- Graphiques avancés non reconstructibles sans Claude Vision",
        "- Tableaux détectés comme blocs colorés uniquement",
    ]
    if debug and ocr_lines:
        lines += ["", "## Textes OCR (debug)", ""]
        for item in ocr_lines[:40]:
            safe = item["text"].replace("|", "\\|")
            lines.append(f"- `{safe}` @ ({item['x']}, {item['y']})")
    return "\n".join(lines)


# ---------------------------------------------------------------------------
# Main pipeline
# ---------------------------------------------------------------------------

def generate_pptx_from_slide_image(
    input_image_path: str,
    output_pptx_path: str,
    output_preview_path: Optional[str] = None,
    output_report_path: Optional[str] = None,
    debug: bool = False,
) -> None:
    print(f"[1/8] Chargement : {input_image_path}")
    img_cv, img_pil = load_image(input_image_path)
    img_h, img_w = img_cv.shape[:2]

    print(f"[2/8] Ratio ({img_w}×{img_h})")
    ratio, ratio_name, slide_dims = detect_slide_ratio(img_cv)
    print(f"      → {ratio_name}  |  slide {slide_dims[0]}×{slide_dims[1]} EMU")

    print("[3/8] Couleurs dominantes")
    colors = extract_dominant_colors(img_cv)

    print("[4/8] Blocs visuels")
    blocks = detect_visual_blocks(img_cv)
    print(f"      → {len(blocks)} blocs")

    print("[5/8] Lignes")
    h_lines, v_lines = detect_lines(img_cv)
    print(f"      → {len(h_lines)} horizontales, {len(v_lines)} verticales")

    print("[6/8] Zones de texte")
    text_zones = detect_text_zones(img_cv)
    print(f"      → {len(text_zones)} zones")

    print(f"[7/8] OCR : {'actif' if TESSERACT_AVAILABLE else 'désactivé'}")
    raw_ocr = ocr_image(img_pil)
    ocr_lines = group_ocr_into_lines(raw_ocr)
    if TESSERACT_AVAILABLE:
        print(f"      → {len(ocr_lines)} lignes extraites")

    Path(output_pptx_path).parent.mkdir(parents=True, exist_ok=True)

    print("[8/8] Génération PPTX")
    prs = build_pptx(img_cv, slide_dims, colors, blocks, h_lines, v_lines, ocr_lines)
    prs.save(output_pptx_path)
    print(f"      → {output_pptx_path}")

    if debug:
        dbg_path = str(Path(output_pptx_path).with_suffix("")) + "_debug.png"
        save_debug_image(img_cv, blocks, h_lines, v_lines, ocr_lines, dbg_path)
        print(f"      → Debug : {dbg_path}")
        print(f"      → Couleurs : {colors[:5]}")

    preview_ok = False
    if output_preview_path:
        print("Aperçu PNG...")
        preview_ok = export_preview(output_pptx_path, output_preview_path)
        status = output_preview_path if preview_ok else "non disponible (LibreOffice requis)"
        print(f"      → {status}")

    if output_report_path:
        report = generate_report(
            input_image_path, img_w, img_h, ratio_name, colors,
            blocks, text_zones, ocr_lines, TESSERACT_AVAILABLE, preview_ok, debug
        )
        Path(output_report_path).parent.mkdir(parents=True, exist_ok=True)
        Path(output_report_path).write_text(report, encoding="utf-8")
        print(f"      → Rapport : {output_report_path}")


# ---------------------------------------------------------------------------
# CLI
# ---------------------------------------------------------------------------

def main() -> None:
    parser = argparse.ArgumentParser(
        description="Convertit une image de slide en PowerPoint éditable."
    )
    parser.add_argument("--input", required=True, help="Image source (PNG/JPG/JPEG)")
    parser.add_argument("--output", required=True, help="Fichier PPTX de sortie")
    parser.add_argument("--preview", default=None, help="Aperçu PNG (optionnel)")
    parser.add_argument("--report", default=None, help="Rapport Markdown (optionnel)")
    parser.add_argument("--debug", action="store_true",
                        help="Image annotée + détails dans le rapport")
    args = parser.parse_args()

    stem = Path(args.output).stem
    out_dir = Path(args.output).parent
    preview = args.preview or str(out_dir / f"{stem}_preview.png")
    report = args.report or str(out_dir / f"{stem}_report.md")

    generate_pptx_from_slide_image(
        input_image_path=args.input,
        output_pptx_path=args.output,
        output_preview_path=preview,
        output_report_path=report,
        debug=args.debug,
    )
    print("\nTerminé.")


if __name__ == "__main__":
    main()
