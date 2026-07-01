#!/usr/bin/env python3
"""
Macro-planning Gantt — modifier la section TACHES pour mettre a jour
pip install python-pptx
Usage : python scripts/gantt_planning.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree


# ════════════════════════════════════════════════════════════
#  DONNÉES — à modifier pour chaque mise à jour
# ════════════════════════════════════════════════════════════

TITRE = "Macro-planning — Chantier 7"

START = (4, 2026)   # (mois, annee) debut du planning
END   = (12, 2026)  # (mois, annee) fin du planning

# debut / fin      : tuple (mois, jour, annee)  ex. (3, 31, 2026) = 31 mars 2026
# fin = None       : jalon ponctuel — aucune barre, seulement le ◆
# jalon_date       : tuple (mois, jour, annee) du ◆ principal — None = pas de diamant
# extra_jalons     : liste de tuples (mois, jour, annee) pour jalons supplementaires
TACHES = [
    {
        "nom":        "Envoie du RFP",
        "jalon":      "J1 — ◆ 15/04/2026",
        "debut":      None, "fin": None,
        "jalon_date": (4, 15, 2026),
    },
    {
        "nom":        "Evaluation des propositions et short présentatitions des providers",
        "jalon":      "J2 — 05/05 → 27/05/2026",
        "debut":      (5, 5, 2026), "fin": (5, 27, 2026),
        "jalon_date": (5, 27, 2026),
    },
    {
        "nom":        "Shortlisting",
        "jalon":      "J3 — ◆ 27/05/2026",
        "debut":      None, "fin": None,
        "jalon_date": (5, 27, 2026),
    },
    {
        "nom":        "Démonstrations des provider",
        "jalon":      "J4 — 02/06 → 04/06/2026",
        "debut":      (6, 2, 2026), "fin": (6, 4, 2026),
        "jalon_date": (6, 4, 2026),
    },
    {
        "nom":        "Sélection du providers",
        "jalon":      "J5 — 02/06 → 12/06/2026",
        "debut":      (6, 2, 2026), "fin": (6, 12, 2026),
        "jalon_date": (6, 12, 2026),
    },
    {
        "nom":        "Signature des lOI et lotissement des chantiers",
        "jalon":      "J6 — date TBD",
        "debut":      (6, 15, 2026), "fin": (12, 31, 2026),
        "jalon_date": None,
    },
    {
        "nom":        "Phase d'implémentation",
        "jalon":      "J7 — date TBD",
        "debut":      (7, 1, 2026), "fin": (12, 31, 2026),
        "jalon_date": None,
    },
]

OUTPUT = "macro_planning_v24.pptx"


# ════════════════════════════════════════════════════════════
#  CHARTE GRAPHIQUE
# ════════════════════════════════════════════════════════════

C_HDR   = RGBColor(0x4E, 0x42, 0x7A)
C_MBG   = RGBColor(0xBA, 0xB6, 0xCC)
C_MTX   = RGBColor(0x4E, 0x42, 0x7A)
C_ODD   = RGBColor(0xD4, 0xD2, 0xE6)
C_EVEN  = RGBColor(0xEC, 0xD8, 0xE4)
C_BAR   = RGBColor(0x68, 0x70, 0xAC)
C_DIA   = RGBColor(0x72, 0x9E, 0x80)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_DARK  = RGBColor(0x20, 0x20, 0x30)
C_GRAY  = RGBColor(0x64, 0x5C, 0x84)

DAYS_IN_MONTH = [31, 28, 31, 30, 31, 30, 31, 31, 30, 31, 30, 31]
MONTH_LABELS  = ["JAN", "FEV", "MAR", "AVR", "MAI", "JUIN",
                 "JUIL", "AOU", "SEP", "OCT", "NOV", "DEC"]


# ════════════════════════════════════════════════════════════
#  GÉOMÉTRIE (pouces)
# ════════════════════════════════════════════════════════════

SL_W, SL_H = 13.33, 7.50

T_L  = 0.25
T_T  = 0.50
T_W  = 12.83
LC_W = 3.00

TL_W = T_W - LC_W   # largeur de la zone timeline

RH_Q = 0.44   # hauteur rangee trimestres
RH_M = 0.32   # hauteur rangee mois
RH_T = 0.63   # hauteur rangee tache

BAR_H = 0.24
BAR_R = 0.22
DIA_S = 0.15
DIA_R = 0.74


# ════════════════════════════════════════════════════════════
#  HELPERS
# ════════════════════════════════════════════════════════════

def get_months_list():
    """Liste des (mois, annee) de START a END inclus."""
    months = []
    m, y = START
    em, ey = END
    while (y, m) <= (ey, em):
        months.append((m, y))
        m += 1
        if m > 12:
            m = 1
            y += 1
    return months


def date_pos(month, day, year):
    """Position fractionnelle (0.0 a N_MONTHS) depuis le debut du planning."""
    sm, sy = START
    return (year - sy) * 12 + (month - sm) + day / DAYS_IN_MONTH[month - 1]


def get_quarters(months_list):
    """Retourne liste de (label, nb_mois) pour l'en-tete trimestriel."""
    groups = []
    current_key = None
    count = 0
    for m, y in months_list:
        q = (m - 1) // 3 + 1
        key = f"T{q} {y}"
        if key != current_key:
            if current_key is not None:
                groups.append((current_key, count))
            current_key = key
            count = 1
        else:
            count += 1
    if current_key:
        groups.append((current_key, count))
    return groups


def _shape(slide, shape_id, l, t, w, h, fill):
    s = slide.shapes.add_shape(shape_id, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.fill.background()
    sp_pr = s._element.find(qn('p:spPr'))
    if sp_pr is not None and sp_pr.find(qn('a:effectLst')) is None:
        etree.SubElement(sp_pr, qn('a:effectLst'))
    return s


def rect(slide, l, t, w, h, fill):
    return _shape(slide, 1, l, t, w, h, fill)


def rounded_rect(slide, l, t, w, h, fill):
    return _shape(slide, 5, l, t, w, h, fill)


def diamond(slide, cx, cy, size, fill):
    h = size / 2
    return _shape(slide, 4, cx - h, cy - h, size, size, fill)


def label(slide, txt, l, t, w, h, size,
          bold=False, italic=False, color=C_WHITE,
          align=PP_ALIGN.CENTER, wrap=False, font_name=None):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = txt
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    if font_name:
        r.font.name = font_name
    return tb


# ════════════════════════════════════════════════════════════
#  GENERATION
# ════════════════════════════════════════════════════════════

def build():
    months_list = get_months_list()
    n_months    = len(months_list)
    M_W         = TL_W / n_months
    quarters    = get_quarters(months_list)

    prs = Presentation()
    prs.slide_width  = Inches(SL_W)
    prs.slide_height = Inches(SL_H)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # En-tete colonne ETAPES
    rect(slide, T_L, T_T, LC_W, RH_Q + RH_M, C_HDR)
    txt_h = 0.22
    txt_y = T_T + (RH_Q + RH_M - txt_h) / 2
    label(slide, "ETAPES", T_L + 0.15, txt_y, LC_W - 0.20, txt_h, 11,
          bold=True, align=PP_ALIGN.LEFT)

    # Trimestres (largeurs variables)
    x_offset = 0.0
    for q_lbl, q_cnt in quarters:
        q_w = q_cnt * M_W
        x = T_L + LC_W + x_offset
        rect(slide, x, T_T, q_w, RH_Q, C_HDR)
        label(slide, q_lbl, x, T_T, q_w, RH_Q, 8, bold=True)
        x_offset += q_w

    # Mois
    y_m = T_T + RH_Q
    for i, (m, y) in enumerate(months_list):
        x = T_L + LC_W + i * M_W
        rect(slide, x, y_m, M_W, RH_M, C_MBG)
        lbl = MONTH_LABELS[m - 1]
        if m == 1 or i == 0:
            lbl = f"{lbl} {y}"
        label(slide, lbl, x, y_m, M_W, RH_M, 6, color=C_MTX)

    # Passe 1 : fonds des lignes de taches
    for i in range(len(TACHES)):
        y_t = T_T + RH_Q + RH_M + i * RH_T
        bg  = C_ODD if i % 2 == 0 else C_EVEN
        rect(slide, T_L, y_t, T_W, RH_T, bg)

    # Passe 2 : filets blancs
    S = 0.008
    total_h = RH_Q + RH_M + len(TACHES) * RH_T
    rect(slide, T_L + LC_W - S / 2, T_T, S, total_h, C_WHITE)
    for i in range(1, n_months):
        x = T_L + LC_W + i * M_W - S / 2
        rect(slide, x, T_T + RH_Q, S, RH_M + len(TACHES) * RH_T, C_WHITE)
    x_offset = 0.0
    for q_lbl, q_cnt in quarters[:-1]:
        x_offset += q_cnt * M_W
        rect(slide, T_L + LC_W + x_offset - S / 2, T_T, S, RH_Q, C_WHITE)
    rect(slide, T_L, T_T + RH_Q + RH_M - S / 2, T_W, S, C_WHITE)

    # Passe 3 : barres, diamants et textes
    for i, t in enumerate(TACHES):
        y_t = T_T + RH_Q + RH_M + i * RH_T

        label(slide, t["nom"],
              T_L + 0.10, y_t + 0.04,
              LC_W - 0.15, RH_T * 0.47,
              8, bold=True, color=C_DARK, align=PP_ALIGN.LEFT, wrap=True,
              font_name="Calibri")

        label(slide, t["jalon"],
              T_L + 0.10, y_t + RH_T * 0.52,
              LC_W - 0.15, RH_T * 0.44,
              6.5, italic=True, color=C_GRAY, align=PP_ALIGN.LEFT, wrap=True,
              font_name="Calibri")

        if t["debut"] is not None and t["fin"] is not None:
            dm, dj, dy = t["debut"]
            fm, fj, fy = t["fin"]
            bx = T_L + LC_W + date_pos(dm, dj, dy) * M_W
            bw = (date_pos(fm, fj, fy) - date_pos(dm, dj, dy)) * M_W
            by = y_t + RH_T * BAR_R
            rounded_rect(slide, bx, by, max(bw, 0.05), BAR_H, C_BAR)

        all_jalons = []
        if t["jalon_date"] is not None:
            all_jalons.append(t["jalon_date"])
        all_jalons.extend(t.get("extra_jalons", []))
        dcy = y_t + RH_T * DIA_R
        for jm, jj, jy in all_jalons:
            dcx = T_L + LC_W + date_pos(jm, jj, jy) * M_W
            diamond(slide, dcx, dcy, DIA_S, C_DIA)

    # Titre sous le tableau
    table_bottom = T_T + RH_Q + RH_M + len(TACHES) * RH_T
    label(slide, TITRE,
          T_L, table_bottom + 0.20, T_W, 0.50,
          13, bold=True, color=C_DARK)

    prs.save(OUTPUT)
    print(f"OK  {OUTPUT} genere ({len(TACHES)} taches, {n_months} mois)")


if __name__ == "__main__":
    build()
