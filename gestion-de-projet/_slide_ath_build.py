#!/usr/bin/env python3
"""Reconstruction: Processus de sélection de la solution cible — v4"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

NAVY  = RGBColor(0x1B, 0x2A, 0x4A)
GOLD  = RGBColor(0xC4, 0x8C, 0x28)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BG    = RGBColor(0xFA, 0xFA, 0xF8)
LGREY = RGBColor(0xCC, 0xD2, 0xDA)
MGREY = RGBColor(0x8A, 0x9B, 0xB5)
PINK  = RGBColor(0xBF, 0x6E, 0x68)
STRIP = RGBColor(0xEF, 0xF2, 0xF5)

FONT_TITLE = "Calibri Light"
FONT_BODY  = "Calibri"


def T(sl, text, x, y, w, h, size, bold=False, italic=False,
      color=NAVY, align=PP_ALIGN.LEFT, wrap=True, font=None):
    font = font or FONT_BODY
    tb = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    first = True
    for line in text.split('\n'):
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color
        r.font.name = font


def S(sl, shape_id, x, y, w, h, fill=None, lc=None, lw=Pt(0)):
    """Generic shape helper."""
    s = sl.shapes.add_shape(shape_id, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if lc:
        s.line.color.rgb = lc; s.line.width = lw
    else:
        s.line.width = 0
    return s


def R(sl, x, y, w, h, fill=None, lc=None, lw=Pt(0)):
    return S(sl, 1, x, y, w, h, fill, lc, lw)


def E(sl, x, y, d, fill=None, lc=GOLD, lw=Pt(1.8)):
    return S(sl, 9, x, y, d, d, fill, lc, lw)


def L(sl, x1, y, x2, color=LGREY, lw=Pt(0.5)):
    c = sl.shapes.add_connector(1, Inches(x1), Inches(y), Inches(x2), Inches(y))
    c.line.color.rgb = color; c.line.width = lw


def chevron(sl, x, y, color=NAVY):
    """Right-pointing chevron as text — renders identically in all viewers."""
    T(sl, "›", x, y, 0.22, 0.36, 18, bold=True, color=color, align=PP_ALIGN.CENTER)


def bullseye(sl, cx, cy, r=0.18):
    """Concentric circles target icon."""
    E(sl, cx - r, cy - r, 2 * r, fill=None, lc=GOLD, lw=Pt(1.2))
    rm = r * 0.58
    E(sl, cx - rm, cy - rm, 2 * rm, fill=GOLD, lc=None)
    rc = r * 0.25
    E(sl, cx - rc, cy - rc, 2 * rc, fill=WHITE, lc=None)


def star(sl, x, y, w=0.22, h=0.22, color=GOLD):
    """5-point star (autoshape 94)."""
    S(sl, 94, x, y, w, h, fill=color)


def cmark(sl, cx, cy, style, r=0.090):
    """Checkmark circle indicator. style: 'full'|'partial'|'empty'."""
    x0, y0, d = cx - r, cy - r, 2 * r
    if style == "full":
        E(sl, x0, y0, d, fill=NAVY, lc=None, lw=Pt(0))
        T(sl, "✓", x0 - 0.005, y0 - 0.01, d + 0.01, d + 0.02, 8,
          bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    elif style == "partial":
        E(sl, x0, y0, d, fill=LGREY, lc=MGREY, lw=Pt(0.4))
        T(sl, "◐", x0, y0 - 0.01, d, d + 0.02, 9,
          color=RGBColor(0x4A, 0x6A, 0x9A), align=PP_ALIGN.CENTER)
    else:
        E(sl, x0, y0, d, fill=None, lc=LGREY, lw=Pt(0.75))


def person_icon(sl, cx, cy, size=0.22, color=NAVY):
    """Simplified person: circle head + wide oval body."""
    hd = size * 0.38
    E(sl, cx - hd / 2, cy - size * 0.05, hd, fill=color, lc=None)
    bw, bh = size * 0.80, size * 0.45
    S(sl, 9, cx - bw / 2, cy + size * 0.30, bw, bh, fill=color)


def two_persons(sl, cx, cy, size=0.26, color=NAVY):
    """Two overlapping person silhouettes."""
    person_icon(sl, cx - size * 0.22, cy, size * 0.85, color)
    person_icon(sl, cx + size * 0.22, cy, size * 0.85, color)


def build():
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.50)
    sl = prs.slides.add_slide(prs.slide_layouts[6])

    # Background + top accent
    bg = sl.shapes.add_shape(1, 0, 0, Inches(13.33), Inches(7.50))
    bg.fill.solid(); bg.fill.fore_color.rgb = BG; bg.line.width = 0
    L(sl, 0.0, 0.02, 13.33, color=NAVY, lw=Pt(2.5))

    # ── TITLE ──────────────────────────────────────────────────────────────
    T(sl, "Processus de sélection de la solution cible",
      0.25, 0.12, 12.3, 0.72, 28, bold=True, font=FONT_TITLE)
    T(sl, "Démarche d'évaluation des éditeurs et aide à la décision",
      0.25, 0.82, 11.0, 0.36, 13, italic=True)
    L(sl, 0.25, 1.28, 13.08)

    # ── PROCESS COLUMNS ────────────────────────────────────────────────────
    COLS = [
        ("1", "CADRAGE", [
            "Finalisation du besoin\net du contenu du RFP",
            "Identification de 10 éditeurs\ncibles",
            "Envoi du RFP et calendrier\nassocié",
        ]),
        ("2", "ÉVALUATION", [
            "Définition des critères d'évaluation\n(fonctionnel, IT, coût, références...)",
            "Pondération par catégorie",
            "Formalisation de la matrice\nd'évaluation",
            "Gestion des questions/réponses\néditeurs",
        ]),
        ("3", "SÉLECTION", [
            "Scoring des réponses\n(fonctionnel, financier, planning...)",
            "Shortlist de 2–3 éditeurs",
            "Organisation des soutenances /\ndémos / POC",
            "Consolidation des évaluations",
        ]),
        ("4", "DÉCISION & DÉPLOIEMENT", [
            "Recommandation finale",
            "Construction du macro–planning\n(POC, migration, budget)",
            "Validation en comité\nde décision",
        ]),
    ]
    CX = [0.25, 3.42, 6.59, 9.76]
    CW = 2.95
    CD = 0.52

    for i, (num, title, bullets) in enumerate(COLS):
        cx = CX[i]
        # Gold circle + number
        ci = cx + (CW - CD) / 2
        E(sl, ci, 1.35, CD)
        T(sl, num, ci + 0.01, 1.37, CD - 0.02, CD - 0.06, 16, color=GOLD,
          align=PP_ALIGN.CENTER, font=FONT_BODY)
        # Column title + gold underline
        T(sl, title, cx, 2.00, CW, 0.35, 8.5, bold=True, align=PP_ALIGN.CENTER)
        L(sl, cx + 0.45, 2.37, cx + CW - 0.45, color=GOLD, lw=Pt(1.2))
        # Bullets
        for k, b in enumerate(bullets):
            T(sl, "▪  " + b, cx + 0.07, 2.47 + k * 0.37, CW - 0.12, 0.36, 7.5)
        # Chevron between columns (centered with circle at y=1.35+0.26=1.61)
        if i < 3:
            chevron(sl, cx + CW + 0.01, 1.44, color=NAVY)

    L(sl, 0.25, 4.10, 13.08)  # SEP 2

    # ── RÔLE DU CABINET DE CONSEIL ─────────────────────────────────────────
    R(sl, 0.25, 4.12, 12.83, 0.57, fill=STRIP)
    # Two-person icon
    two_persons(sl, 0.52, 4.41, size=0.26, color=NAVY)
    T(sl, "RÔLE DU CABINET\nDE CONSEIL", 0.72, 4.15, 1.48, 0.52, 7.5, bold=True)

    # Role descriptions with small icons
    ROLES = [
        ("Structurer la démarche\net cadrer le besoin",    2.30, 2.75, "target"),
        ("Construire le cadre\nd'évaluation",               5.10, 1.90, "list"),
        ("Soutien opérationnel\nchef de projet IT +\nchef de projet Métier", 7.05, 3.00, "person"),
        ("Consolider les analyses\net formuler la recommandation", 10.10, 2.95, "chart"),
    ]
    for txt, rx, rw, icon_type in ROLES:
        # Small icon approximation
        ix = rx - 0.25
        iy = 4.32
        if icon_type == "target":
            bullseye(sl, ix + 0.12, iy + 0.12, r=0.10)
        elif icon_type == "list":
            for li in range(3):
                R(sl, ix, iy + li * 0.07, 0.18, 0.04, fill=GOLD)
        elif icon_type == "person":
            person_icon(sl, ix + 0.10, iy + 0.05, size=0.20, color=NAVY)
        elif icon_type == "chart":
            for bi, bh in enumerate([0.09, 0.18, 0.12]):
                R(sl, ix + bi * 0.065, iy + 0.18 - bh, 0.05, bh, fill=GOLD)
        T(sl, txt, rx, 4.15, rw, 0.54, 7.5)

    L(sl, 0.25, 4.72, 13.08)  # SEP 3

    # ── BOTTOM: CRITÈRES D'ÉVALUATION (gauche) ─────────────────────────────
    T(sl, "CRITÈRES D'ÉVALUATION", 0.25, 4.76, 6.0, 0.24, 7.5, bold=True)
    L(sl, 0.25, 5.00, 6.0, color=GOLD, lw=Pt(1.5))

    HX = [3.20, 4.10, 5.00]
    for j, h in enumerate(["PRESTATAIRE A", "PRESTATAIRE B", "PRESTATAIRE C"]):
        T(sl, h, HX[j], 5.02, 0.88, 0.21, 5.5, bold=True, align=PP_ALIGN.CENTER)

    CRIT = [
        ("Couverture fonctionnelle",              "full",    "full",    "partial"),
        ("Ergonomie & expérience utilisateur",    "full",    "partial", "full"),
        ("Architecture & intégration SI",         "full",    "full",    "partial"),
        ("Références & expérience marché",        "full",    "full",    "partial"),
        ("Coût total de possession",              "partial", "full",    "partial"),
        ("Planning & capacité de déploiement",    "full",    "partial", "partial"),
        ("Accompagnement au changement",          "full",    "partial", "partial"),
        ("Sécurité & conformité réglementaire",   "full",    "full",    "partial"),
        ("Roadmap produit & innovation",          "full",    "partial", "full"),
    ]
    for k, (crit, a, b, c) in enumerate(CRIT):
        ry = 5.25 + k * 0.205
        # Row separator line between rows
        if k > 0:
            L(sl, 0.25, ry, 5.88, color=LGREY, lw=Pt(0.3))
        T(sl, crit, 0.48, ry + 0.01, 2.68, 0.19, 7.0)
        row_cy = ry + 0.1025
        for j, style in enumerate([a, b, c]):
            cmark(sl, HX[j] + 0.44, row_cy, style)
    # Bottom border of table
    L(sl, 0.25, 5.25 + 9 * 0.205, 5.88, color=LGREY, lw=Pt(0.3))

    T(sl, "●  Répond au critère    ◐  Répond partiellement    ○  Ne répond pas",
      0.25, 7.02, 6.0, 0.20, 6.0, color=MGREY)

    # ── BOTTOM: PODIUM (centre) ─────────────────────────────────────────────
    T(sl, "POSITIONNEMENT DES SOLUTIONS ÉVALUÉES",
      6.30, 4.76, 4.80, 0.24, 7.5, bold=True)
    L(sl, 6.30, 5.00, 9.65, color=GOLD, lw=Pt(1.5))

    BASE = 7.05
    BARS = [("B", 6.50, 1.10, MGREY), ("A", 7.28, 1.68, GOLD), ("C", 8.06, 0.90, PINK)]
    for label, bx, bh, bc in BARS:
        by = BASE - bh
        R(sl, bx, by, 0.68, bh, fill=bc)
        T(sl, label, bx, by + bh / 2 - 0.18, 0.68, 0.36, 17, bold=True,
          color=WHITE, align=PP_ALIGN.CENTER)
    # Laurel star above A bar
    star(sl, 7.28 + 0.68 / 2 - 0.11, BASE - 1.68 - 0.26, w=0.22, h=0.22, color=GOLD)
    L(sl, 6.38, BASE, 8.90, color=NAVY, lw=Pt(0.75))

    # ── BOTTOM: RECOMMANDATION (droite) ────────────────────────────────────
    R(sl, 10.25, 4.74, 2.80, 2.52, fill=WHITE, lc=LGREY, lw=Pt(0.75))
    # Bullseye icon centered above "RECOMMANDATION"
    bullseye(sl, 11.65, 4.90, r=0.10)
    T(sl, "RECOMMANDATION", 10.30, 5.04, 2.70, 0.22, 8, bold=True,
      color=NAVY, align=PP_ALIGN.CENTER)
    L(sl, 10.50, 5.28, 12.87, color=GOLD, lw=Pt(1.0))
    # Trophy star centered, above "SOLUTION RECOMMANDÉE"
    star(sl, 11.54, 5.33, w=0.22, h=0.22, color=GOLD)
    T(sl, "SOLUTION\nRECOMMANDÉE", 10.30, 5.58, 2.70, 0.44, 9.5, bold=True,
      color=NAVY, align=PP_ALIGN.CENTER)
    T(sl, "PRESTATAIRE A", 10.30, 6.06, 2.70, 0.30, 13, bold=True,
      color=GOLD, align=PP_ALIGN.CENTER)
    T(sl,
      "Le Prestataire A est recommandé, étant le seul à couvrir l'ensemble "
      "des besoins prioritaires identifiés lors du cadrage, avec une adéquation "
      "optimale aux enjeux fonctionnels, techniques et organisationnels.",
      10.28, 6.40, 2.74, 0.82, 7.0, wrap=True)

    # ── FOOTER ─────────────────────────────────────────────────────────────
    L(sl, 0.25, 7.30, 13.08)
    T(sl, "Soutien opérationnel chef de projet IT  +  Chef de projet Métier  +  Métier",
      0.25, 7.33, 12.83, 0.18, 8, italic=True, align=PP_ALIGN.CENTER)

    return prs


if __name__ == "__main__":
    out = Path(r"C:\Users\arnau\OneDrive\Bureau\slide_ath_v4.pptx")
    build().save(str(out))
    print(f"Sauvegardé : {out}")
